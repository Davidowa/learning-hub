"""Slide layouts for the Academia de Cómputo deck template.

One class, one method per slide archetype. Every method drops a finished slide
into the deck and returns it, so a lesson file reads as a list of slides rather
than a pile of coordinates.

Text is measured with the real font files before it is placed, so a heading that
runs long shrinks to fit instead of spilling off the canvas.
"""
from __future__ import annotations

import os

from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import tokens as K
from .highlight import highlight

HERE = os.path.dirname(os.path.abspath(__file__))
BRAND = os.path.join(os.path.dirname(HERE), 'brand')

_FONT_FILES = {
    (K.SANS, False): 'arial.ttf',  (K.SANS, True): 'arialbd.ttf',
    (K.SERIF, False): 'georgia.ttf', (K.SERIF, True): 'georgiab.ttf',
    (K.MONO, False): 'cour.ttf',   (K.MONO, True): 'courbd.ttf',
}
_font_cache: dict = {}


def _pil(name: str, bold: bool, size_pt: float):
    key = (name, bold, round(size_pt, 1))
    if key not in _font_cache:
        path = os.path.join(os.environ.get('WINDIR', r'C:\Windows'), 'Fonts',
                            _FONT_FILES.get((name, bold), 'arial.ttf'))
        try:
            _font_cache[key] = ImageFont.truetype(path, max(4, int(round(size_pt * 4))))
        except OSError:
            _font_cache[key] = None
    return _font_cache[key]


def measure(text: str, name: str, size_pt: float, bold: bool = False) -> float:
    """Width of ``text`` in inches at ``size_pt``."""
    f = _pil(name, bold, size_pt)
    if f is None:                      # no font file — fall back to a ratio
        return len(text) * size_pt * 0.0072
    return f.getlength(text) / 4 / 72.0


def wrap_lines(text: str, name: str, size_pt: float, width_in: float, bold=False) -> int:
    """How many lines ``text`` needs inside ``width_in``."""
    usable = width_in - 0.07           # the 25400 EMU body insets
    lines, cur = 1, ''
    for word in text.split():
        trial = f'{cur} {word}'.strip()
        if cur and measure(trial, name, size_pt, bold) > usable:
            lines, cur = lines + 1, word
        else:
            cur = trial
    return lines


def fit_size(text: str, name: str, size_pt: float, width_in: float, max_lines: int,
             bold=False, floor=18.0) -> float:
    """Largest size at or below ``size_pt`` that keeps ``text`` within ``max_lines``."""
    size = size_pt
    while size > floor and wrap_lines(text, name, size, width_in, bold) > max_lines:
        size -= 0.75
    return size


def line_h(size_pt: float, spacing: float = 1.0) -> float:
    """Height of one rendered line, in inches. PowerPoint uses 1.2 x em."""
    return size_pt * 1.2 * spacing / 72.0


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


class Deck:
    """A presentation in the Academia de Cómputo house style."""

    def __init__(self, footer: str, lang: str = 'es', course: str = '',
                 blocks: int = 3, notes: bool = True, language: str = 'python'):
        self.C = K.palette(language)
        self.prs = Presentation()
        self.prs.slide_width = Inches(K.SLIDE_W)
        self.prs.slide_height = Inches(K.SLIDE_H)
        self.footer = footer
        self.lang = lang
        self.course = course
        self.blocks = blocks          # progress pips in the footer
        self.block = 0                # which pip is lit
        self.notes = notes
        self._blank = self.prs.slide_layouts[6]

    # ───────────────────────────────────────────── primitives

    def _slide(self, bg: str):
        s = self.prs.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(bg)
        return s

    def _box(self, s, x, y, w, h, shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.background()
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def _rect(self, s, x, y, w, h, fill=None, line=None, lw=0.75, radius=None):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if radius:
            sp.adjustments[0] = radius
        if fill:
            sp.fill.solid()
            sp.fill.fore_color.rgb = rgb(fill)
        else:
            sp.fill.background()
        if line:
            sp.line.color.rgb = rgb(line)
            sp.line.width = Pt(lw)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        sp.text_frame.word_wrap = True
        return sp

    def _rule(self, s, x, y, w, color=None, h=0.012):
        color = color or self.C.BORDER
        return self._rect(s, x, y, w, h, fill=color)

    def _text(self, s, x, y, w, h, runs, size=K.T.body, font=K.SANS, color=None,
              bold=False, spacing=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              space_after=0.0):
        """``runs`` is a string, or a list of (text, font, size, colour, bold)."""
        color = color or self.C.INK
        sp = self._box(s, x, y, w, h)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(25400)
        p = tf.paragraphs[0]
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        if space_after:
            p.space_after = Pt(space_after)
        items = [(runs,)] if isinstance(runs, (str, int, float)) else list(runs)
        for item in items:
            item = tuple(item) + (None,) * (5 - len(item))
            t, fnt, sz, col, bd = item[:5]
            r = p.add_run()
            r.text = str(t)
            r.font.name = fnt or font
            r.font.size = Pt(sz or size)
            r.font.bold = bool(bold if bd is None else bd)
            r.font.color.rgb = rgb(col or color)
            r.font._rPr.set('lang', 'es-MX' if self.lang == 'es' else 'en-US')
        return sp

    def _pic(self, s, path, x, y, w, h, alt=''):
        p = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        p._element._nvXxPr.cNvPr.set('descr', alt)
        return p

    def _dot(self, s, x, y, d, fill):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def _shadow(self, sp, blur=28, dist=8, alpha=16000):
        """Carbon-style soft drop shadow, straight into the shape XML."""
        spPr = sp._element.spPr
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        spPr.append(parse_xml(
            f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:outerShdw blurRad="{blur * 12700}" dist="{dist * 12700}" dir="5400000" '
            f'rotWithShape="0"><a:srgbClr val="0B1B3A"><a:alpha val="{alpha}"/></a:srgbClr>'
            f'</a:outerShdw></a:effectLst>'))

    # ───────────────────────────────────────────── slide chrome

    def _eyebrow(self, s, text, color=None):
        color = color or self.C.BLUE
        self._text(s, K.M, K.EYEBROW_Y, K.SLIDE_W - 2 * K.M, 0.32, text.upper(),
                   size=K.T.eyebrow, font=K.MONO, color=color)

    def _title(self, s, text, width=K.TITLE_W, color=None):
        """Place the slide title; return how far a wrapped title pushes content down."""
        color = color or self.C.INK
        size = fit_size(text, K.SANS, K.T.title, width, 2, bold=True, floor=30)
        lines = wrap_lines(text, K.SANS, size, width, bold=True)
        lh = size * 1.2 * K.LS.title / 72.0
        self._text(s, K.M, K.TITLE_Y, width, lh * lines + 0.2, text,
                   size=size, color=color, bold=True, spacing=K.LS.title)
        return max(0.0, (lines - 1) * lh)

    def _head(self, s, eyebrow, title, eyebrow_color=None):
        self._eyebrow(s, eyebrow, eyebrow_color)
        return self._title(s, title)

    def _footer(self, s, page):
        self._rule(s, K.M, K.FOOTER_RULE_Y, K.CONTENT_W, self.C.BORDER, 0.012)
        logo = os.path.join(BRAND, 'duck-512.png')
        if os.path.exists(logo):
            self._pic(s, logo, *K.FOOTER_LOGO, alt='')
        self._text(s, 1.65, K.FOOTER_Y, 9.5, 0.32, self.footer.upper(),
                   size=K.T.caption, font=K.MONO, color=self.C.MUTED)
        for i in range(self.blocks):
            x = 17.72 - (self.blocks - 1 - i) * 0.79
            self._rect(s, x, 10.26, 0.54, 0.05,
                       fill=self.C.BLUE_FILL if i < self.block else self.C.BORDER)
        self._text(s, K.PAGENUM_X, K.FOOTER_Y, 0.38, 0.32, f'{page:02d}',
                   size=K.T.caption, font=K.MONO, color=self.C.INK)

    def _page(self, s, notes=None):
        self._footer(s, len(self.prs.slides))
        if notes and self.notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    def _dark_page(self, s, notes=None):
        self._text(s, K.PAGENUM_X, K.FOOTER_Y, 0.38, 0.32, f'{len(self.prs.slides):02d}',
                   size=K.T.caption, font=K.MONO, color=self.C.ON_NAVY_DIM)
        if notes and self.notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    # ───────────────────────────────────────────── 01 cover

    def cover(self, kicker, unit, title, subtitle, meta=(), notes=None):
        s = self._slide(self.C.NAVY)
        logo = os.path.join(BRAND, 'duck-on-dark-512.png')
        if os.path.exists(logo):
            self._pic(s, logo, 1.04, 1.00, 0.62, 0.62,
                      alt='Rubber duck, the mark of the Academia de Cómputo')
        self._text(s, 1.88, 1.16, 8.0, 0.34, kicker.upper(), size=K.T.label,
                   font=K.MONO, color=self.C.ON_NAVY_LINK)
        self._text(s, K.M, 4.31, 15.58, 0.32, unit.upper(), size=K.T.eyebrow,
                   font=K.MONO, color=self.C.ACCENT)
        size = fit_size(title, K.SERIF, K.T.cover_title, 14.59, 3, floor=54)
        n = wrap_lines(title, K.SERIF, size, 14.59)
        lh = size * 1.2 * K.LS.cover_title / 72.0
        self._text(s, K.M, 7.10 - n * lh, 14.59, n * lh + 0.2, title, size=size,
                   font=K.SERIF, color=self.C.ON_NAVY, spacing=K.LS.cover_title)
        self._text(s, K.M, 7.58, 12.4, 1.11, subtitle, size=K.T.cover_sub,
                   color=self.C.ON_NAVY_SOFT, spacing=K.LS.lead)
        self._rule(s, K.M, 9.31, K.CONTENT_W, self.C.PAPER, 0.012)
        for i, (label, value) in enumerate(meta):
            x = K.M + i * 4.15
            self._text(s, x, 9.70, 3.92, 0.32, label.upper(), size=K.T.eyebrow,
                       font=K.MONO, color=self.C.ON_NAVY_DIM)
            self._text(s, x, 10.06, 3.92, 0.40, value, size=K.T.value, color=self.C.ON_NAVY)
        return self._dark_page(s, notes)

    # ───────────────────────────────────────────── 02 agenda / four cards

    def agenda(self, eyebrow, title, cards, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        top = K.CARD_Y + dy
        cards = cards[:4]
        tw, dw = K.CARD_W - 0.68, K.CARD_W - 0.80

        # Measure every card first, then give all four the tallest one's height so
        # the row stays a grid even when one blurb runs long.
        laid = []
        for c in cards:
            ts = fit_size(c['title'], K.SANS, K.T.card_title, tw, 2, True)
            tl = wrap_lines(c['title'], K.SANS, ts, tw, True)
            dl = wrap_lines(c['desc'], K.SANS, K.T.card_body, dw)
            laid.append((ts, tl, dl))
        th = [line_h(ts, K.LS.card_title) for ts, _, _ in laid]
        dh = line_h(K.T.card_body, K.LS.body)
        body = max(1.08 + tl * th[i] + 0.14 + dl * dh for i, (_, tl, dl) in enumerate(laid))
        h = max(K.CARD_H, body + (0.86 if any(c.get('meta') for c in cards) else 0.34))

        for i, (c, (ts, tl, dl)) in enumerate(zip(cards, laid)):
            x = K.M + i * K.CARD_STEP
            card = self._rect(s, x, top, K.CARD_W, h, fill=self.C.WHITE, line=self.C.BORDER)
            self._shadow(card, blur=18, dist=4, alpha=7000)
            self._text(s, x + K.PAD, top + 0.39, K.CARD_W - 0.6, 0.56,
                       c.get('num', f'{i + 1:02d}'), size=K.T.card_num, font=K.MONO,
                       color=self.C.BLUE)
            self._text(s, x + K.PAD, top + 1.08, tw, tl * th[i] + 0.2, c['title'],
                       size=ts, bold=True, spacing=K.LS.card_title)
            self._text(s, x + K.PAD, top + 1.08 + tl * th[i] + 0.14, dw,
                       dl * dh + 0.2, c['desc'], size=K.T.card_body, color=self.C.MUTED,
                       spacing=K.LS.body)
            if c.get('meta'):
                self._text(s, x + K.PAD, top + h - 0.62, K.CARD_W - 0.6, 0.34,
                           c['meta'], size=K.T.label, font=K.MONO, color=self.C.INK)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 03 objectives

    def objectives(self, eyebrow, title, items, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        for i, it in enumerate(items[:6]):
            col, row = i % 2, i // 2
            ax = K.M + col * 8.50
            y = 2.61 + dy + row * 2.15
            self._text(s, ax, y, 0.30, 0.40, '→', size=K.T.value, font=K.MONO,
                       color=self.C.ACCENT_LO)
            hs = fit_size(it['head'], K.SANS, K.T.head, 7.9, 2, True, floor=20)
            self._text(s, ax + 0.46, y - 0.04, 7.9, 0.95, it['head'], size=hs,
                       bold=True, spacing=K.LS.head)
            self._text(s, ax + 0.46, y + 0.51, 7.9, 0.9, it['sub'], size=K.T.sub,
                       color=self.C.MUTED, spacing=K.LS.sub)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 04 roadmap

    def roadmap(self, eyebrow, title, phases, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        step = K.CONTENT_W / max(1, len(phases))
        w = step - 0.34

        # Measure every phase before drawing anything. The active phase gets a card,
        # and a fixed height leaves a three-line description hanging off the bottom
        # of it, on the paper. The card has to cover the tallest phase, not its own.
        lh = line_h(K.T.label, K.LS.desc)
        laid = []
        for p in phases:
            ts = fit_size(p['title'], K.SANS, K.T.card_title, w - 0.4, 2, True, floor=19)
            tl = wrap_lines(p['title'], K.SANS, ts, w - 0.4, True)
            dl = wrap_lines(p['desc'], K.SANS, K.T.label, w - 0.5)
            laid.append((ts, tl, dl))
        deepest = max(4.05 + (0.42 if tl > 1 else 0) + dl * lh for _, tl, dl in laid)
        card_h = max(2.33, deepest + 0.30 - 2.65)

        for i, (p, (ts, tl, dl)) in enumerate(zip(phases, laid)):
            x = K.M + i * step
            state = p.get('state', 'todo')
            if state == 'now':
                act = self._rect(s, x, 2.65 + dy, w, card_h, fill=self.C.WHITE,
                                 line=self.C.BORDER)
                self._shadow(act, blur=18, dist=4, alpha=8000)
                self._rect(s, x, 2.65 + dy, w, 0.06, fill=self.C.ACCENT_LO)
                lab, tcol = self.C.ACCENT_LO, self.C.INK
                tx = x + 0.33
            else:
                self._rect(s, x, 2.65 + dy, w, 0.06,
                           fill=self.C.BLUE_FILL if state == 'done' else self.C.BORDER)
                lab = self.C.BLUE if state == 'done' else self.C.MUTED
                tcol = self.C.INK if state == 'done' else self.C.MUTED
                tx = x if i == 0 else x + 0.33
            self._text(s, tx, 3.05 + dy, w - 0.3, 0.34, p['label'].upper(),
                       size=K.T.label, font=K.MONO, color=lab)
            self._text(s, tx, 3.48 + dy, w - 0.4, 0.9, p['title'], size=ts, bold=True,
                       color=tcol, spacing=K.LS.card_title)
            self._text(s, tx, 4.05 + dy + (0.42 if tl > 1 else 0), w - 0.5,
                       dl * lh + 0.2, p['desc'], size=K.T.label,
                       color=self.C.MUTED, spacing=K.LS.desc)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 05 block divider

    def divider(self, eyebrow, title, lede, block=None, notes=None):
        if block is not None:
            self.block = block
        s = self._slide(self.C.NAVY)
        self._text(s, K.M, 3.59, 15.0, 0.32, eyebrow.upper(), size=K.T.eyebrow,
                   font=K.MONO, color=self.C.ACCENT)
        size = fit_size(title, K.SERIF, K.T.divider_title, 15.02, 2, floor=48)
        self._text(s, K.M, 4.21, 15.02, 2.23, title, size=size, font=K.SERIF,
                   color=self.C.ON_NAVY, spacing=K.LS.divider_title)
        self._text(s, K.M, 6.77, 12.0, 1.2, lede, size=K.T.divider_lede,
                   color=self.C.ON_NAVY_SOFT, spacing=K.LS.lead)
        return self._dark_page(s, notes)

    # ───────────────────────────────────────────── 06 concept

    def concept(self, eyebrow, title, lead, panel=None, steps=None, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        top = 2.57 + dy
        self._text(s, K.M, top, 8.6, 1.9, lead, size=K.T.lead, spacing=K.LS.lead)
        y = top + 0.44 * wrap_lines(lead if isinstance(lead, str) else
                                    ''.join(r[0] for r in lead),
                                    K.SANS, K.T.lead, 8.6) + 0.55
        if panel:
            self._rule(s, K.M, y, 8.58)
            self._text(s, K.M, y + 0.30, 9.44, 0.34, panel['label'].upper(),
                       size=K.T.label, font=K.MONO, color=self.C.MUTED)
            self._text(s, K.M, y + 0.79, 9.44, 0.5, panel['title'], size=K.T.value,
                       spacing=K.LS.body)
            if panel.get('note'):
                self._text(s, K.M, y + 1.40, 9.44, 0.9, panel['note'], size=K.T.sub,
                           color=self.C.MUTED, spacing=K.LS.loose)
        if steps:
            h = 0.95 + 0.81 * len(steps) + (0.95 if steps and steps[-1].get('loop') else 0)
            card = self._rect(s, 10.38, top, 8.58, max(3.4, h), fill=self.C.WHITE, line=self.C.BORDER)
            self._shadow(card, blur=22, dist=5, alpha=10000)
            for i, st in enumerate(steps):
                yy = top + 0.47 + i * 0.81
                accent = self.C.ACCENT_LO if st.get('accent') else self.C.BLUE_FILL
                self._dot(s, 10.84, yy, 0.58, accent)
                self._text(s, 10.80, yy, 0.67, 0.62, str(i + 1), size=K.T.eyebrow,
                           font=K.MONO, color=self.C.WHITE, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)
                self._text(s, 11.64, yy + 0.11, 6.9, 0.5, st['text'], size=K.T.value)
            if steps and steps[-1].get('loop'):
                self._text(s, 10.84, top + 0.47 + len(steps) * 0.81 + 0.14, 7.88, 0.8,
                           steps[-1]['loop'], size=K.T.label, font=K.MONO,
                           color=self.C.MUTED, spacing=K.LS.desc)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 07 pipeline diagram

    def diagram(self, eyebrow, title, blocks, note=None, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        n = len(blocks)
        gap = 1.25
        w = (K.CONTENT_W - gap * (n - 1)) / n
        for i, b in enumerate(blocks):
            x = K.M + i * (w + gap)
            dark = b.get('dark')
            blk = self._rect(s, x, 2.65 + dy, w, 2.48, fill=self.C.NAVY if dark else self.C.WHITE,
                       line=None if dark else self.C.BORDER, radius=0.05)
            self._text(s, x + 0.32, 3.04 + dy, w - 0.5, 0.34, b['label'].upper(),
                       size=K.T.label, font=K.MONO,
                       color=self.C.ACCENT if dark else self.C.BLUE)
            self._text(s, x + 0.32, 3.49 + dy, w - 0.5, 0.5, b['title'], size=K.T.block_title,
                       bold=True, color=self.C.ON_NAVY if dark else self.C.INK,
                       spacing=K.LS.card_title)
            self._shadow(blk, blur=18, dist=4, alpha=7000)
            self._text(s, x + 0.32, 4.02 + dy, w - 0.7, 1.0, b['desc'], size=K.T.label,
                       color=self.C.ON_NAVY_SOFT if dark else self.C.MUTED, spacing=K.LS.desc)
            if i < n - 1:
                self._text(s, x + w + 0.34, 3.63 + dy, 0.6, 0.6, '→', size=K.T.arrow,
                           color=self.C.MUTED)
        if note:
            self._text(s, K.M, 5.45 + dy, 15.0, 1.0, note, size=K.T.note, color=self.C.MUTED,
                       spacing=K.LS.note)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 08–11 code card

    def code(self, eyebrow, title, filename, source, lang='python', theme='dark',
             annotations=(), width=K.CODE_W, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        self._code_card(s, K.CODE_X, K.CODE_Y + dy, width, filename, source, lang, theme)
        y = 2.57 + dy
        for a in annotations[:4]:
            n = wrap_lines(a['text'], K.SANS, K.T.note, K.ANNOT_W)
            self._text(s, K.ANNOT_X, y, K.ANNOT_W, 0.34, a['label'].upper(),
                       size=K.T.label, font=K.MONO,
                       color=self.C.ACCENT_LO if a.get('accent') else self.C.BLUE)
            self._text(s, K.ANNOT_X, y + 0.41, K.ANNOT_W, n * 0.42 + 0.2, a['text'],
                       size=K.T.note, spacing=K.LS.note)
            y += 0.41 + n * 0.42 + 0.36
        return self._page(s, notes)

    def _code_card(self, s, x, y, w, filename, source, lang='python', theme='dark',
                   maxlines=K.CODE_MAX_LINES, size=K.T.code, bottom=K.FOOTER_RULE_Y - 0.36):
        """Draw a carbon-style code card and return its bottom edge.

        The type shrinks (never below the 18 pt projection floor) until the longest
        line fits the width and every line fits above ``bottom``.
        """
        th = self.C.DARK_CODE if theme == 'dark' else self.C.LIGHT_CODE
        all_lines = highlight(source, lang)
        lines = all_lines[:maxlines]
        if len(all_lines) > maxlines:
            print(f'  ! {filename}: {len(all_lines)} lines of code, card shows {maxlines}')

        # A blank line separates, it does not carry code, so it gets a half step.
        BLANK = 0.5
        weights = [BLANK if not any(t.strip() for t, _ in ln) else 1.0 for ln in lines]
        rows = sum(weights) or 1

        widest = max((sum(len(t) for t, _ in ln) for ln in lines), default=1)
        by_width = (w - 0.95) / (widest * 0.6 / 72.0)          # Courier is 0.6 em wide
        room = bottom - y - K.CODE_HEADER_H - 0.36 - 0.22
        by_height = room / rows / (K.CODE_LINE_STEP / K.T.code)
        size = min(size, by_width, by_height)
        # quantise to fixed steps: three code slides in a row should not each land
        # on a different fractional size
        size = max(K.T.eyebrow,
                   K.T.eyebrow + K.CODE_STEP * int((size - K.T.eyebrow) / K.CODE_STEP))
        step = K.CODE_LINE_STEP * (size / K.T.code)
        h = K.CODE_HEADER_H + 0.36 + step * rows + 0.22
        if y + h > bottom + 0.30:
            print(f'  ! {filename}: {len(lines)} lines do not fit above the footer at '
                  f'the 18 pt floor; split the example across two slides')

        card = self._rect(s, x, y, w, h, fill=th['bg'],
                          line=None if theme == 'dark' else self.C.BORDER, radius=0.035)
        self._shadow(card, blur=30, dist=9, alpha=14000 if theme == 'dark' else 8000)
        if th['header']:
            self._rect(s, x + 0.01, y + 0.01, w - 0.02, K.CODE_HEADER_H, fill=th['header'])
        self._rule(s, x + 0.01, y + K.CODE_HEADER_H, w - 0.02, th['rule'], 0.012)
        # carbon-style window dot + file name
        self._dot(s, x + 0.34, y + 0.28, 0.14, th['dot'])
        self._text(s, x + 0.62, y + 0.21, w - 1.0, 0.32, filename, size=K.T.caption,
                   font=K.MONO, color=th['chrome'])
        cursor = y + K.CODE_HEADER_H + 0.36
        for spans, weight in zip(lines, weights):
            if weight == 1.0:
                runs = [(t, K.MONO, size, th.get(role, th['text']), False)
                        for t, role in spans]
                self._text(s, x + 0.34, cursor, w - 0.6, 0.5, runs, spacing=K.LS.code)
            cursor += step * weight
        return y + h

    # ───────────────────────────────────────────── illustration

    def figure(self, eyebrow, title, image, caption=None, alt=None, side=None,
               notes=None):
        """A generated doodle diagram, centred, with an optional note beside it.

        ``image`` is a path relative to the lesson folder or to ppts/. A missing
        file is reported and the slide still builds, so a deck never blocks on an
        illustration that has not been generated yet.
        """
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        top = 2.61 + dy
        room_w = (11.4 if side else K.CONTENT_W)
        # the caption sits under the picture, so take its height out of the box
        # before scaling; otherwise a full-width figure pushes it into the footer
        cap_lines = wrap_lines(caption, K.SANS, K.T.note, room_w) if caption else 0
        cap_h = (0.22 + cap_lines * line_h(K.T.note, K.LS.note)) if caption else 0
        room_h = K.FOOTER_RULE_Y - 0.30 - cap_h - top

        path = image if os.path.isabs(image) else os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))), image)
        if not os.path.exists(path):
            print(f'  ! figure missing: {image}')
            self._rect(s, K.M, top, room_w, room_h, fill=self.C.PANEL,
                       line=self.C.BORDER, radius=0.04)
            self._text(s, K.M + 0.5, top + room_h / 2 - 0.3, room_w - 1.0, 0.6,
                       f'[ {image} ]', size=K.T.body, font=K.MONO, color=self.C.MUTED,
                       align=PP_ALIGN.CENTER)
        else:
            from PIL import Image as _Im
            with _Im.open(path) as im:
                iw, ih = im.size
            scale = min(room_w / iw, room_h / ih)
            w, h = iw * scale, ih * scale
            self._pic(s, path, K.M + (room_w - w) / 2, top, w, h,
                      alt=alt or caption or title)
            if caption:
                self._text(s, K.M, top + h + 0.22, room_w, 0.6, caption,
                           size=K.T.note, color=self.C.MUTED, spacing=K.LS.note,
                           align=PP_ALIGN.CENTER if not side else PP_ALIGN.LEFT)
        if side:
            self._text(s, 12.90, top + 0.10, 6.06, 0.34, side['label'].upper(),
                       size=K.T.label, font=K.MONO, color=self.C.BLUE)
            self._text(s, 12.90, top + 0.56, 6.06, 3.0, side['text'], size=K.T.body,
                       spacing=K.LS.body)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 10 code + console output

    def code_output(self, eyebrow, title, filename, source, output, lang='python',
                    output_label='Salida', notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        top = K.CODE_Y + dy
        self._code_card(s, K.CODE_X, top, K.CODE_W, filename, source, lang, 'dark')
        ox, ow = 12.03, 6.93
        out_lines = output.rstrip('\n').split('\n')[:10]
        h = 1.06 + K.CODE_LINE_STEP * max(1, len(out_lines)) + 0.22
        panel = self._rect(s, ox, top, ow, h, fill=self.C.WHITE, line=self.C.BORDER, radius=0.035)
        self._shadow(panel, blur=20, dist=5, alpha=7000)
        self._text(s, ox + 0.34, top + 0.24, ow - 0.6, 0.34, output_label.upper(),
                   size=K.T.label, font=K.MONO, color=self.C.BLUE)
        self._rule(s, ox + 0.01, top + K.CODE_HEADER_H, ow - 0.02, self.C.BORDER, 0.012)
        for i, line in enumerate(out_lines):
            self._text(s, ox + 0.34, top + K.CODE_HEADER_H + 0.36 +
                       i * K.CODE_LINE_STEP, ow - 0.6, 0.5, line or ' ',
                       size=K.T.code, font=K.MONO, color=self.C.INK, spacing=K.LS.code)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 12 comparison table

    def table(self, eyebrow, title, headers, rows, accent_row=None, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        n = len(headers)
        widths = [K.CONTENT_W * (0.19 if i == 0 else 0.81 / (n - 1)) for i in range(n)]
        xs, acc = [], K.M
        for w in widths:
            xs.append(acc)
            acc += w
        y = 2.65 + dy
        for i, htxt in enumerate(headers):
            self._text(s, xs[i], y, widths[i] - 0.3, 0.34, htxt.upper(),
                       size=K.T.label, font=K.MONO, color=self.C.MUTED)
        y += 0.52
        self._rule(s, K.M, y, K.CONTENT_W, self.C.INK, 0.016)
        y += 0.22
        for r, row in enumerate(rows[:8]):
            if len(row) != n:
                raise ValueError(f'table row {r + 1} has {len(row)} cells but there '
                                 f'are {n} headers: {row!r}. A comma inside a YAML '
                                 f'flow list splits the cell; quote it.')
            hs = max(wrap_lines(str(c), K.MONO if i else K.SANS, K.T.body,
                                widths[i] - 0.35) for i, c in enumerate(row))
            # 0.38, not 0.42: seven single-line rows have to clear the footer
            # rule, and the convention is to show a group of operators whole
            rh = 0.38 + 0.40 * hs
            if r % 2 == 1:
                self._rect(s, K.M - 0.12, y - 0.10, K.CONTENT_W + 0.24, rh, fill=self.C.PANEL)
            if accent_row == r:
                self._rect(s, K.M - 0.12, y - 0.10, 0.06, rh, fill=self.C.ACCENT_LO)
            for i, cell in enumerate(row):
                mono = i > 0
                self._text(s, xs[i], y, widths[i] - 0.3, rh, str(cell),
                           size=K.T.sub if mono else K.T.body,
                           font=K.MONO if mono else K.SANS,
                           color=self.C.INK if i == 0 or accent_row == r else self.C.MUTED,
                           bold=(i == 0), spacing=K.LS.body)
            y += rh + 0.14
        return self._page(s, notes)

    # ───────────────────────────────────────────── 13 numbered method

    def steps(self, eyebrow, title, items, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        n = min(len(items), 5)
        gap = 0.42
        w = (K.CONTENT_W - gap * (n - 1)) / n
        for i, it in enumerate(items[:n]):
            x = K.M + i * (w + gap)
            accent = it.get('accent')
            self._rect(s, x, 2.65 + dy, w, 0.06,
                       fill=self.C.ACCENT_LO if accent else self.C.BLUE_FILL)
            self._text(s, x, 3.03 + dy, w, 0.34, it.get('label', f'{i + 1:02d}').upper(),
                       size=K.T.label, font=K.MONO,
                       color=self.C.ACCENT_LO if accent else self.C.BLUE)
            ts = fit_size(it['title'], K.SANS, K.T.card_title, w - 0.2, 3, True, floor=19)
            tl = wrap_lines(it['title'], K.SANS, ts, w - 0.2, True)
            self._text(s, x, 3.46 + dy, w - 0.2, tl * 0.44 + 0.2, it['title'], size=ts,
                       bold=True, spacing=K.LS.card_title)
            self._text(s, x, 3.46 + dy + tl * 0.44 + 0.18, w - 0.3, 1.8, it['desc'],
                       size=K.T.card_body, color=self.C.MUTED, spacing=K.LS.body)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 14 big number

    def stat(self, eyebrow, value, text, note=None, notes=None):
        s = self._slide(self.C.NAVY)
        self._text(s, K.M, 2.60, 15.0, 0.32, eyebrow.upper(), size=K.T.eyebrow,
                   font=K.MONO, color=self.C.ACCENT)
        self._text(s, K.M, 3.25, 14.0, 2.4, value, size=K.T.stat, font=K.SERIF,
                   color=self.C.ON_NAVY, spacing=0.9)
        self._text(s, K.M, 6.10, 13.0, 1.6, text, size=K.T.divider_lede,
                   color=self.C.ON_NAVY, spacing=K.LS.lead)
        if note:
            self._text(s, K.M, 8.00, 12.0, 1.0, note, size=K.T.body,
                       color=self.C.ON_NAVY_SOFT, spacing=K.LS.body)
        return self._dark_page(s, notes)

    # ───────────────────────────────────────────── 16 error catalogue

    def pitfalls(self, eyebrow, title, items, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        for i, it in enumerate(items[:4]):
            col, row = i % 2, i // 2
            x = K.M + col * 9.29
            y = 2.61 + dy + row * 2.45
            self._rect(s, x, y, 0.06, 1.55, fill=self.C.ACCENT_LO)
            self._text(s, x + 0.36, y, 8.3, 0.34,
                       it.get('label', f'ERROR {i + 1:02d}').upper(),
                       size=K.T.label, font=K.MONO, color=self.C.ACCENT_LO)
            self._text(s, x + 0.36, y + 0.44, 8.3, 0.5, it['title'],
                       size=K.T.card_title, bold=True, spacing=K.LS.card_title)
            self._text(s, x + 0.36, y + 1.00, 8.3, 0.9, it['desc'], size=K.T.sub,
                       color=self.C.MUTED, spacing=K.LS.sub)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 17 method / pull quote

    def method(self, eyebrow, title, statement, points=(), notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        top = 2.70 + dy
        self._text(s, K.M, top, 16.5, 1.9, statement, size=K.T.statement, font=K.SERIF,
                   spacing=1.06)
        y = top + 0.62 * wrap_lines(statement, K.SERIF, K.T.statement, 16.5) + 0.80
        self._rule(s, K.M, y, K.CONTENT_W)
        n = max(1, len(points))
        w = (K.CONTENT_W - 0.8 * (n - 1)) / n
        for i, p in enumerate(points):
            x = K.M + i * (w + 0.8)
            self._text(s, x, y + 0.40, w, 0.34, p['label'].upper(), size=K.T.label,
                       font=K.MONO, color=self.C.BLUE)
            self._text(s, x, y + 0.86, w, 1.4, p['text'], size=K.T.body,
                       color=self.C.MUTED, spacing=K.LS.body)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 18 lab brief

    def lab(self, eyebrow, title, brief, blocks=(), notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title, eyebrow_color=self.C.ACCENT_LO)
        bs = fit_size(brief, K.SANS, K.T.lead, 16.5, 3, floor=21)
        bl = wrap_lines(brief, K.SANS, bs, 16.5)
        bh = 0.62 + bl * 0.52
        card = self._rect(s, K.M, 2.60 + dy, K.CONTENT_W, bh, fill=self.C.WHITE,
                          line=self.C.BORDER, radius=0.04)
        self._shadow(card, blur=20, dist=5, alpha=8000)
        self._text(s, K.M + 0.5, 2.90 + dy, 16.5, bl * 0.52 + 0.2, brief, size=bs,
                   spacing=K.LS.lead)
        n = max(1, len(blocks))
        w = (K.CONTENT_W - 0.7 * (n - 1)) / n
        for i, b in enumerate(blocks):
            x = K.M + i * (w + 0.7)
            self._text(s, x, 2.60 + dy + bh + 0.50, w, 0.34, b['label'].upper(),
                       size=K.T.label, font=K.MONO, color=self.C.MUTED)
            self._text(s, x, 2.60 + dy + bh + 0.96, w, 1.6, b['text'], size=K.T.body,
                       spacing=K.LS.body)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 19 quiz

    def quiz(self, eyebrow, question, filename=None, source=None, lang='python',
             options=(), notes=None):
        s = self._slide(self.C.PAPER)
        self._eyebrow(s, eyebrow, self.C.ACCENT_LO)
        size = fit_size(question, K.SANS, K.T.title, 16.5, 2, True, floor=32)
        self._text(s, K.M, K.TITLE_Y, 16.5, 1.5, question, size=size, bold=True,
                   spacing=K.LS.title)
        tl = wrap_lines(question, K.SANS, size, 16.5, True)
        top = 3.10 + (tl - 1) * 0.78
        if source:
            self._code_card(s, K.M, top, 9.6, filename or 'quiz.py', source, lang, 'dark')
        for i, opt in enumerate(options[:4]):
            y = top + 0.20 + i * 1.10
            x = 11.4 if source else K.M
            self._rect(s, x, y, 0.62, 0.62, fill=self.C.WHITE, line=self.C.BORDER, radius=0.5)
            self._text(s, x, y, 0.62, 0.62, chr(65 + i), size=K.T.value, font=K.MONO,
                       color=self.C.BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._text(s, x + 0.90, y + 0.08, 6.6, 0.9, opt, size=K.T.body,
                       spacing=K.LS.body)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 20 trace table

    def trace(self, eyebrow, title, headers, rows, verdict=None, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        n = len(headers)
        w = K.CONTENT_W / n
        y = 2.70 + dy
        for i, h in enumerate(headers):
            self._text(s, K.M + i * w, y, w - 0.3, 0.34, h.upper(), size=K.T.label,
                       font=K.MONO, color=self.C.MUTED)
        y += 0.52
        self._rule(s, K.M, y, K.CONTENT_W, self.C.INK, 0.016)
        y += 0.20
        for r, row in enumerate(rows[:8]):
            hs = max(wrap_lines(str(c), K.MONO, K.T.sub, w - 0.3) for c in row)
            rh = 0.30 + 0.40 * hs
            if r % 2 == 1:
                self._rect(s, K.M - 0.12, y - 0.08, K.CONTENT_W + 0.24, rh, fill=self.C.PANEL)
            for i, cell in enumerate(row):
                self._text(s, K.M + i * w, y, w - 0.3, rh, str(cell), size=K.T.sub,
                           font=K.MONO, color=self.C.INK if i == 0 else self.C.MUTED,
                           spacing=K.LS.body)
            y += rh + 0.14
        if verdict:
            self._rect(s, K.M, y + 0.25, 0.06, 0.75, fill=self.C.ACCENT_LO)
            self._text(s, K.M + 0.36, y + 0.28, 15.0, 0.8, verdict, size=K.T.lead,
                       spacing=K.LS.lead)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 21 before / after

    def compare(self, eyebrow, title, before, after, note=None, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        bottom = 0.0
        for i, (side, theme) in enumerate(((before, 'light'), (after, 'dark'))):
            x = K.M + i * 9.29
            self._text(s, x, 2.49 + dy, 8.6, 0.34, side['label'].upper(), size=K.T.label,
                       font=K.MONO, color=self.C.MUTED if i == 0 else self.C.ACCENT_LO)
            bottom = max(bottom, self._code_card(
                s, x, 2.92 + dy, 8.6, side.get('file', ''), side['code'],
                side.get('lang', 'python'), theme, size=19.5))
        if note:
            self._rect(s, K.M, bottom + 0.55, 0.06, 0.80, fill=self.C.ACCENT_LO)
            self._text(s, K.M + 0.36, bottom + 0.55, 16.4, 1.0, note, size=K.T.note,
                       color=self.C.MUTED, spacing=K.LS.note)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 22 cost / tiers

    def tiers(self, eyebrow, title, items, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        for i, it in enumerate(items[:4]):
            y = 2.80 + dy + i * 1.62
            self._text(s, K.M, y, 3.0, 0.6, it['key'], size=K.T.card_num, font=K.MONO,
                       color=self.C.ACCENT_LO if it.get('accent') else self.C.BLUE)
            self._text(s, 4.30, y + 0.05, 5.2, 0.5, it['title'], size=K.T.card_title,
                       bold=True, spacing=K.LS.card_title)
            self._text(s, 9.80, y + 0.05, 9.0, 1.0, it['desc'], size=K.T.body,
                       color=self.C.MUTED, spacing=K.LS.body)
            if i < len(items[:4]) - 1:
                self._rule(s, K.M, y + 1.26, K.CONTENT_W)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 23 quotation

    def quote(self, text, author, source=None, notes=None):
        s = self._slide(self.C.NAVY)
        size = fit_size(f'«{text}»', K.SERIF, K.T.quote, 16.4, 5, floor=32)
        self._text(s, K.M, 3.00, 16.4, 3.6, f'«{text}»', size=size, font=K.SERIF,
                   color=self.C.ON_NAVY, spacing=1.14)
        self._text(s, K.M, 7.30, 12.0, 0.5, author, size=K.T.value, color=self.C.ACCENT)
        if source:
            self._text(s, K.M, 7.85, 12.0, 0.5, source, size=K.T.eyebrow, font=K.MONO,
                       color=self.C.ON_NAVY_DIM)
        return self._dark_page(s, notes)

    # ───────────────────────────────────────────── 25 / 27 numbered list

    def takeaways(self, eyebrow, title, items, notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        for i, it in enumerate(items[:4]):
            y = 2.70 + dy + i * 1.65
            self._text(s, K.M, y, 1.0, 0.6, f'{i + 1:02d}', size=K.T.card_num,
                       font=K.MONO, color=self.C.BLUE)
            self._text(s, 2.20, y + 0.02, 8.0, 0.6, it['title'], size=K.T.head,
                       bold=True, spacing=K.LS.head)
            self._text(s, 2.20, y + 0.62, 15.0, 0.8, it['desc'], size=K.T.sub,
                       color=self.C.MUTED, spacing=K.LS.sub)
            if i < len(items[:4]) - 1:
                self._rule(s, K.M, y + 1.35, K.CONTENT_W)
        return self._page(s, notes)

    # ───────────────────────────────────────────── 26 homework

    def homework(self, eyebrow, title, brief, meta=(), rubric=(), notes=None):
        s = self._slide(self.C.PAPER)
        dy = self._head(s, eyebrow, title)
        bs = fit_size(brief, K.SANS, K.T.lead, 16.8, 3, floor=21)
        self._text(s, K.M, 2.57 + dy, 16.8, 1.5, brief, size=bs, spacing=K.LS.lead)
        y = 2.57 + dy + 0.46 * wrap_lines(brief, K.SANS, bs, 16.8) + 0.45
        n = max(1, len(meta))
        w = (K.CONTENT_W - 0.7 * (n - 1)) / n
        for i, (label, value) in enumerate(meta):
            x = K.M + i * (w + 0.7)
            self._text(s, x, y, w, 0.34, label.upper(), size=K.T.label, font=K.MONO,
                       color=self.C.MUTED)
            self._text(s, x, y + 0.44, w, 0.9, value, size=K.T.body, spacing=K.LS.body)
        if rubric:
            y2 = y + 1.60
            self._rule(s, K.M, y2, K.CONTENT_W, self.C.INK, 0.016)
            y2 += 0.30
            for row in rubric[:4]:
                if len(row) != 3:
                    raise ValueError(f'rubric row needs 3 cells (label, weight, '
                                     f'description), got {len(row)}: {row!r}. '
                                     f'A comma inside a YAML flow list splits the '
                                     f'cell; quote it.')
                label, weight, desc = row
                self._text(s, K.M, y2, 5.0, 0.5, label, size=K.T.body, bold=True)
                self._text(s, 6.20, y2, 10.6, 0.5, desc, size=K.T.sub, color=self.C.MUTED)
                self._text(s, 16.9, y2, 2.06, 0.5, weight, size=K.T.body, font=K.MONO,
                           color=self.C.BLUE, align=PP_ALIGN.RIGHT)
                y2 += 0.70
        return self._page(s, notes)

    # ───────────────────────────────────────────── 28 closing

    def closing(self, title, subtitle, meta=(), notes=None):
        s = self._slide(self.C.NAVY)
        self._text(s, K.M, 3.60, 14.0, 1.9, title, size=K.T.closing_title, font=K.SERIF,
                   color=self.C.ON_NAVY, spacing=0.92)
        self._text(s, K.M, 5.90, 13.0, 1.0, subtitle, size=K.T.divider_lede,
                   color=self.C.ON_NAVY_SOFT, spacing=K.LS.lead)
        self._rule(s, K.M, 7.60, K.CONTENT_W, self.C.ON_NAVY_DIM, 0.012)
        for i, (label, value) in enumerate(meta):
            x = K.M + i * 4.60
            self._text(s, x, 8.00, 4.3, 0.32, label.upper(), size=K.T.eyebrow,
                       font=K.MONO, color=self.C.ON_NAVY_DIM)
            self._text(s, x, 8.40, 4.3, 0.5, value, size=K.T.value, color=self.C.ON_NAVY)
        return self._dark_page(s, notes)

    # ───────────────────────────────────────────── output

    def save(self, path):
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
        return path
