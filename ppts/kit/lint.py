"""Check a built deck for text that spills past the safe area.

Measures every run with the real font metrics, wraps it inside its own box, and
reports anything whose rendered bottom crosses the footer rule or whose right
edge crosses the margin.

    python -m kit.lint out/w03.es.pptx
    python -m kit.lint out
"""
from __future__ import annotations

import argparse
import glob
import os

from pptx import Presentation

from . import tokens as K
from .deck import line_h, measure, wrap_lines

SAFE_RIGHT = K.M + K.CONTENT_W + 0.10      # 19.06 in
FOOTER_TOP = K.FOOTER_RULE_Y - 0.06        # content must end above the rule
CHROME_TOP = K.SLIDE_H - 0.12              # the footer band itself


def check(path: str) -> list[str]:
    prs = Presentation(path)
    problems = []
    for n, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            x = sh.left / 914400
            y = sh.top / 914400
            w = sh.width / 914400
            right = x + w

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                # Zebra bands and rules bleed past the margin on purpose, so an
                # empty shape is not measured sideways. It still has to respect
                # the footer rule downwards: a code card whose type has bottomed
                # out at the 18 pt floor keeps growing, and the card runs under
                # the rule while its last line of type stops just above it. Every
                # text check therefore passes and the deck reads clean, which is
                # how twenty-seven overflowing cards once shipped.
                #
                # Only a panel is checked. A rule is 0.01 in tall, the progress
                # dashes are 0.05, and a full-bleed background starts at the top
                # edge, so none of them can trip this.
                bottom = y + sh.height / 914400
                if sh.height / 914400 > 0.5 and y > 0.10 and bottom > K.FOOTER_RULE_Y + 0.02:
                    problems.append(f'{n:02d}  panel runs to {bottom:.2f}" and crosses '
                                    f'the footer rule at {K.FOOTER_RULE_Y}"')
                continue

            if right > SAFE_RIGHT + 0.01:
                problems.append(f'{n:02d}  text box runs to {right:.2f}" '
                                f'(margin ends at {SAFE_RIGHT:.2f}")')
            para = sh.text_frame.paragraphs[0]
            runs = para.runs
            if not runs:
                continue
            size = max((r.font.size.pt for r in runs if r.font.size), default=18)
            name = runs[0].font.name or K.SANS
            bold = bool(runs[0].font.bold)
            text = ''.join(r.text for r in runs)
            mult = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
            lines = wrap_lines(text, name, size, w, bold)
            bottom = y + lines * line_h(size, mult)

            limit = CHROME_TOP if y >= FOOTER_TOP else FOOTER_TOP
            if bottom > limit + 0.05:
                snippet = text[:52].replace('\n', ' ')
                problems.append(f'{n:02d}  text ends at {bottom:.2f}" '
                                f'(limit {limit:.2f}"): "{snippet}"')

            # A code line and an output line each get their own box, at the code
            # line spacing. Either one wrapping means the tail is sitting on the
            # slide background, which the checks above never see because the box
            # itself stays inside the safe area.
            #
            # Measure the whole string rather than asking wrap_lines: that helper
            # splits on words, which collapses runs of spaces, so an aligned
            # trailing comment measures short and a line that really overflows
            # comes back as fitting. PowerPoint keeps every space.
            if name == K.MONO and abs(mult - K.LS.code) < 1e-6:
                used = measure(text, name, size)
                if used > w - 0.07:
                    snippet = text.strip()[:48]
                    problems.append(f'{n:02d}  code line overflows its card '
                                    f'({used:.2f}" in {w - 0.07:.2f}", '
                                    f'{len(text)} chars): "{snippet}"')

    return problems


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split('\n')[0])
    ap.add_argument('target', help='a .pptx, or a folder of them')
    a = ap.parse_args(argv)
    files = (sorted(glob.glob(os.path.join(a.target, '**', '*.pptx'), recursive=True))
             if os.path.isdir(a.target) else [a.target])
    total = 0
    for f in files:
        found = check(f)
        total += len(found)
        mark = 'clean' if not found else f'{len(found)} issue(s)'
        print(f'{os.path.basename(f):<20} {mark}')
        for p in found:
            print(f'    {p}')
    print(f'\n{total} issue(s) across {len(files)} deck(s)')
    return 1 if total else 0


if __name__ == '__main__':
    raise SystemExit(main())
