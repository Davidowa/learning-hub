"""Report every place the kit shrank type below its nominal size.

The layouts auto-fit: a long heading drops a point at a time until it wraps to the
allowed number of lines, and a code card scales to the widest line. This prints
what actually landed on the slides so the shrinking stays visible instead of
silent.

    python -m kit.sizes python
"""
from __future__ import annotations

import argparse
import glob
import os
from collections import Counter

from pptx import Presentation

from . import tokens as K

FLOOR = 18.0                    # the projection floor, ~24 px

# nominal size -> what that size is used for, so a shrink can be named
NOMINAL = {
    K.T.cover_title: 'cover title', K.T.divider_title: 'divider title',
    K.T.title: 'slide title', K.T.quote: 'pull quote', K.T.lead: 'lead paragraph',
    K.T.head: 'objective heading', K.T.card_title: 'card title',
    K.T.code: 'code', K.T.card_body: 'card body',
}


def audit(path):
    prs = Presentation(path)
    sizes = Counter()
    shrunk, below = [], []
    for n, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if not r.font.size:
                        continue
                    pt = r.font.size.pt
                    sizes[pt] += 1
                    if pt < FLOOR:
                        below.append((n, pt, r.text[:40]))
                    # a size that is not one of the tokens means auto-fit moved it
                    elif pt not in vars(K.T).values():
                        shrunk.append((n, pt, r.text[:46]))
    return sizes, shrunk, below


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split('\n')[0])
    ap.add_argument('target')
    a = ap.parse_args(argv)
    files = (sorted(glob.glob(os.path.join(a.target, '**', '*.pptx'), recursive=True))
             if os.path.isdir(a.target) else [a.target])

    total_below = 0
    for f in files:
        sizes, shrunk, below = audit(f)
        print(f'\n{os.path.basename(f)}')
        if not sizes:
            # a deck of nothing but pictures carries no sized runs, and min() on an
            # empty counter used to abort the whole run before the generated decks
            print('  no sized text runs, nothing to check')
            continue
        print(f'  smallest {min(sizes):.2f} pt   largest {max(sizes):.2f} pt   '
              f'{len(sizes)} distinct sizes')
        if shrunk:
            print(f'  auto-fit moved {len(shrunk)} run(s) off a token size:')
            seen = set()
            for n, pt, text in shrunk:
                key = (n, round(pt, 2))
                if key in seen:
                    continue
                seen.add(key)
                print(f'    slide {n:02d}  {pt:5.2f} pt  "{text}"')
        else:
            print('  every run sits on a token size')
        if below:
            total_below += len(below)
            print(f'  !! {len(below)} run(s) below the {FLOOR} pt floor:')
            for n, pt, text in below[:8]:
                print(f'    slide {n:02d}  {pt:5.2f} pt  "{text}"')
    print(f'\n{total_below} run(s) below {FLOOR} pt across {len(files)} deck(s)')
    return 1 if total_below else 0


if __name__ == '__main__':
    raise SystemExit(main())
