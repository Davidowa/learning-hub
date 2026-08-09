"""Rasterise a generated deck to PNG contact sheets.

The kit only ever emits rectangles, rounded rectangles, ovals, pictures and
top-anchored text boxes, so a small renderer covers the whole vocabulary. Handy
for checking a build without opening PowerPoint.

    python -m kit.preview out/com102-w03.es.pptx out/preview --cols 3
"""
from __future__ import annotations

import argparse
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .deck import _FONT_FILES

DPI = 100
INSET = 0.0278


def _font(name, bold, size_pt, px_per_in):
    fname = _FONT_FILES.get((name, bold), 'arial.ttf')
    path = os.path.join(os.environ.get('WINDIR', r'C:\Windows'), 'Fonts', fname)
    px = max(6, int(round(size_pt / 72.0 * px_per_in)))
    try:
        return ImageFont.truetype(path, px)
    except OSError:
        return ImageFont.load_default()


def _hex(color_fmt):
    try:
        if color_fmt and color_fmt.type is not None:
            return '#' + str(color_fmt.rgb)
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _wrap(draw, runs, width_px):
    """Greedy wrap over (text, font, colour) runs. Returns list of line-run-lists."""
    lines, cur, cur_w = [], [], 0
    for text, font, col in runs:
        parts = text.split(' ')
        for i, word in enumerate(parts):
            piece = word if i == 0 else ' ' + word
            if not piece:
                continue
            w = draw.textlength(piece, font=font)
            if cur and cur_w + w > width_px and piece.strip():
                lines.append(cur)
                cur, cur_w = [], 0
                piece = piece.lstrip(' ')
                w = draw.textlength(piece, font=font)
            cur.append((piece, font, col))
            cur_w += w
    if cur:
        lines.append(cur)
    return lines or [[]]


def render_slide(slide, w_in, h_in, scale=DPI, bg='#FFFFFF'):
    W, H = int(w_in * scale), int(h_in * scale)
    img = Image.new('RGB', (W, H), bg)
    try:
        b = _hex(slide.background.fill.fore_color)
        if b:
            img = Image.new('RGB', (W, H), b)
    except (AttributeError, TypeError, ValueError):
        pass
    draw = ImageDraw.Draw(img)

    for sh in slide.shapes:
        x, y = sh.left / 914400 * scale, sh.top / 914400 * scale
        w, h = sh.width / 914400 * scale, sh.height / 914400 * scale
        box = [x, y, x + w, y + h]

        if sh.shape_type == 13:                                   # picture
            try:
                im = Image.open(__import__('io').BytesIO(sh.image.blob)).convert('RGBA')
                im = im.resize((max(1, int(w)), max(1, int(h))), Image.LANCZOS)
                img.paste(im, (int(x), int(y)), im)
            except Exception:
                pass
            continue

        fill = line = None
        try:
            if sh.fill.type is not None and str(sh.fill.type).startswith('SOLID'):
                fill = _hex(sh.fill.fore_color)
        except (AttributeError, TypeError, ValueError):
            pass
        try:
            if str(sh.line.fill.type).startswith('SOLID'):
                line = _hex(sh.line.color)
        except (AttributeError, TypeError, ValueError):
            pass
        name = str(sh.shape_type)
        if fill or line:
            lw = max(1, int((sh.line.width or 0) / 12700 / 72 * scale)) if line else 0
            if 'OVAL' in name:
                draw.ellipse(box, fill=fill, outline=line, width=lw or 1)
            elif 'ROUNDED' in name:
                r = min(w, h) * float(sh.adjustments[0] or 0)
                draw.rounded_rectangle(box, radius=max(1, r), fill=fill,
                                       outline=line, width=lw or 1)
            else:
                draw.rectangle(box, fill=fill, outline=line, width=lw or 1)

        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        tf = sh.text_frame
        pad = INSET * scale
        tx, ty = x + pad, y + pad
        avail = max(10, w - 2 * pad)
        for para in tf.paragraphs:
            runs = []
            for r in para.runs:
                fnt = _font(r.font.name or 'Arial', bool(r.font.bold),
                            r.font.size.pt if r.font.size else 18, scale)
                col = '#000000'
                try:
                    col = '#' + str(r.font.color.rgb)
                except (AttributeError, TypeError, ValueError):
                    pass
                runs.append((r.text, fnt, col))
            if not runs:
                continue
            size_pt = max((r.font.size.pt for r in para.runs if r.font.size), default=18)
            mult = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
            lh = size_pt / 72 * scale * 1.2 * mult
            lines = _wrap(draw, runs, avail)
            total = lh * len(lines)
            oy = ty
            if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
                oy = y + (h - total) / 2
            for ln in lines:
                lw_px = sum(draw.textlength(t, font=f) for t, f, _ in ln)
                ox = tx
                if para.alignment == PP_ALIGN.CENTER:
                    ox = x + (w - lw_px) / 2
                elif para.alignment == PP_ALIGN.RIGHT:
                    ox = x + w - pad - lw_px
                asc = max((f.getmetrics()[0] for _, f, _ in ln), default=lh)
                for t, f, c in ln:
                    draw.text((ox, oy + (lh - asc) * 0.55), t, font=f, fill=c)
                    ox += draw.textlength(t, font=f)
                oy += lh
            ty = oy
    return img


def render(pptx_path, out_prefix, cols=3, scale=DPI, sheet_width=1500):
    prs = Presentation(pptx_path)
    w_in = prs.slide_width / 914400
    h_in = prs.slide_height / 914400
    imgs = [render_slide(s, w_in, h_in, scale) for s in prs.slides]
    os.makedirs(os.path.dirname(os.path.abspath(out_prefix)) or '.', exist_ok=True)

    tw = sheet_width // cols
    th = int(tw * h_in / w_in)
    per = cols * 3
    made = []
    for page in range(0, len(imgs), per):
        chunk = imgs[page:page + per]
        rows = (len(chunk) + cols - 1) // cols
        sheet = Image.new('RGB', (cols * tw + 8, rows * (th + 20) + 8), '#D8DEE9')
        d = ImageDraw.Draw(sheet)
        f = _font('Courier New', False, 9, 96)
        for i, im in enumerate(chunk):
            cx, cy = (i % cols) * tw + 4, (i // cols) * (th + 20) + 4
            sheet.paste(im.resize((tw - 8, th - 8), Image.LANCZOS), (cx, cy))
            d.text((cx + 2, cy + th - 4), f'{page + i + 1:02d}', font=f, fill='#0B1B3A')
        p = f'{out_prefix}-{page // per + 1}.png'
        sheet.save(p)
        made.append(p)
    return made


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx')
    ap.add_argument('out_prefix')
    ap.add_argument('--cols', type=int, default=3)
    ap.add_argument('--scale', type=int, default=DPI)
    a = ap.parse_args()
    for p in render(a.pptx, a.out_prefix, a.cols, a.scale):
        print('wrote', p)
